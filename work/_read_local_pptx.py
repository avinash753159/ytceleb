"""Read notes AND comments from the owner's downloaded copy of the deck.

The Google copy was read through the Slides API, which surfaces speaker notes
and Drive comments. A .pptx edited in PowerPoint can also carry comment parts
that never round-tripped to Drive, so this opens the file directly: notes via
python-pptx, comments by unzipping and reading the raw comment XML.

Anything found here is diffed against manifest/deck_feedback.json so only new
input is reported.
"""
import json
import re
import zipfile
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
LOCAL = Path(r"C:\Users\avina\Downloads"
             r"\MrBeast V8 — picture review (172 slides).pptx")
KNOWN = ROOT / "manifest/deck_feedback.json"
OUT = ROOT / "manifest/deck_feedback_local.json"


def norm(t: str) -> str:
    return re.sub(r"\s+", " ", (t or "")).strip()


def main() -> int:
    if not LOCAL.exists():
        cands = sorted(Path(r"C:\Users\avina\Downloads").glob("*.pptx"))
        print(f"not found: {LOCAL}")
        print("pptx files in Downloads:")
        for c in cands:
            print("   ", c.name)
        return 1

    known = {}
    if KNOWN.exists():
        for n in json.loads(KNOWN.read_text(encoding="utf-8"))["notes"]:
            known[n["slide"]] = norm(n["added_note"])

    # ---- comments, straight out of the package ------------------------
    print("=== COMMENT PARTS IN THE PACKAGE ===")
    found_comments = []
    with zipfile.ZipFile(LOCAL) as z:
        names = [n for n in z.namelist()
                 if "comment" in n.lower() and n.endswith(".xml")]
        if not names:
            print("(none - no comment parts in the file)")
        for n in names:
            xml = z.read(n).decode("utf-8", "ignore")
            texts = re.findall(r"<(?:a:t|p188:txBody[^>]*)>(.*?)</a:t>", xml,
                               re.S)
            body = [norm(re.sub(r"<[^>]+>", " ", t)) for t in texts]
            body = [b for b in body if b]
            if body:
                print(f"--- {n}")
                for b in body:
                    print("   ", b[:300])
                found_comments.append({"part": n, "text": body})

    # ---- speaker notes ------------------------------------------------
    prs = Presentation(str(LOCAL))
    print(f"\n=== NOTES ({len(prs.slides.__iter__.__self__._sldIdLst)} slides) ===")
    new, changed = [], []
    for i, sl in enumerate(prs.slides, 1):
        if not sl.has_notes_slide:
            continue
        txt = norm(sl.notes_slide.notes_text_frame.text)
        if not txt:
            continue
        # strip the generated block so only the owner's words remain
        extra = txt
        for marker in ("why this picture:",):
            if marker in extra:
                extra = extra.split(marker, 1)[1]
        extra = norm(extra)
        if not extra:
            continue
        prev = known.get(i, "")
        if prev and (prev in extra or extra in prev):
            continue
        # is any of this genuinely new?
        if prev:
            changed.append({"slide": i, "was": prev[:200], "now": extra[:600]})
        else:
            new.append({"slide": i, "note": extra[:600]})

    print(f"\nslides whose note is NOT already in deck_feedback.json: "
          f"{len(new)} new, {len(changed)} changed")
    for r in new[:40]:
        print(f"\n--- SLIDE {r['slide']} (new)")
        print("   ", r["note"][:400])
    for r in changed[:40]:
        print(f"\n--- SLIDE {r['slide']} (changed)")
        print("    now:", r["now"][:400])

    OUT.write_text(json.dumps(
        {"comments": found_comments, "new": new, "changed": changed},
        indent=2), encoding="utf-8")
    print(f"\n[OK] {OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
