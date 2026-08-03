#!/usr/bin/env python3
"""Build the slideshow as a deck for review, one slide per shot.

Slide N IS shot N - there is no title slide - so "change slide 47" is
unambiguous. The shot number is also printed on every slide in case anyone
counts differently.

Each slide carries the frame, the programme timecode, what kind of picture it
is, which asset it came from, and the line being spoken over it, because the
question being reviewed is almost always "does this picture belong on this
sentence". The speaker notes carry the plan's stated intent for the segment
and the full asset path, so the slide itself stays clean.

Built as .pptx and uploaded with conversion rather than assembled through the
Slides API: the API's createImage needs a publicly fetchable URL for every
image, which would mean making 172 frames of an unreleased film link-readable.
A converted upload embeds them and shares nothing.
"""

from __future__ import annotations

import json
import os
import subprocess
import sys
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "pipeline"))

SUFFIX = os.environ.get("BUILD_SUFFIX", "V8")
PIECES = ROOT / f"work/picture_{SUFFIX.lower()}"
SHOTS = ROOT / f"manifest/picture_{SUFFIX.lower()}_shots.json"
EDL = ROOT / "manifest/edl_full.json"
STILLS = ROOT / f"work/slides_{SUFFIX.lower()}"
OUT = ROOT / f"final_video/MRBEAST_{SUFFIX}_SLIDESHOW.pptx"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
IMG_X, IMG_Y, IMG_W, IMG_H = Inches(0.667), Inches(0.12), Inches(12.0), Inches(6.75)
INK = RGBColor(0xEA, 0xEA, 0xEA)
DIM = RGBColor(0x9A, 0x9A, 0x9A)
ACCENT = RGBColor(0xE3, 0x12, 0x0B)
GROUND = RGBColor(0x0A, 0x0A, 0x0C)


def tc(t: float) -> str:
    return f"{int(t // 60)}:{t % 60:05.2f}"


def still(s: dict) -> Path:
    """The frame from the MIDDLE of the shot - every shot fades in, so a first
    frame would give a deck of fade-ups."""
    src = PIECES / f"{s['name']}.mp4"
    dest = STILLS / f"{s['name']}.jpg"
    if dest.exists():
        return dest
    subprocess.run(
        ["ffmpeg", "-hide_banner", "-loglevel", "error", "-y",
         "-ss", f"{s['dur']*0.5:.3f}", "-i", str(src), "-frames:v", "1",
         "-vf", "scale=1280:-2", "-q:v", "4", str(dest)],
        check=True, timeout=300)
    return dest


def spoken(seg_by_i: dict, s: dict) -> str:
    span = str(s.get("span", ""))
    try:
        seg = seg_by_i.get(int(float(span)))
    except (TypeError, ValueError):
        seg = None
    if span == "6.5":
        return "[TITLE BREAK - music only, no narration]"
    if span == "-1":
        return "[music lead-in, before the first word]"
    if span == "TAIL":
        return "[music tail, after the last word]"
    if not seg:
        return ""
    if seg["kind"] == "beat":
        return "[beat - music only]"
    return (seg.get("text") or "").strip()


def main() -> int:
    STILLS.mkdir(parents=True, exist_ok=True)
    shots = json.loads(SHOTS.read_text(encoding="utf-8"))
    edl = json.loads(EDL.read_text(encoding="utf-8"))
    seg_by_i = {s["i"]: s for s in edl["segs"]}

    print(f"[deck] {len(shots)} shots -> extracting stills", flush=True)
    with ThreadPoolExecutor(max_workers=8) as ex:
        for f in as_completed([ex.submit(still, s) for s in shots]):
            f.result()

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    for n, s in enumerate(shots, 1):
        sl = prs.slides.add_slide(blank)
        bg = sl.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = GROUND
        bg.line.fill.background()
        bg.shadow.inherit = False

        sl.shapes.add_picture(str(still(s)), IMG_X, IMG_Y, IMG_W, IMG_H)

        # red rule under the frame, then the caption line
        rule = sl.shapes.add_shape(1, IMG_X, Inches(6.97), Inches(0.06),
                                   Inches(0.42))
        rule.fill.solid()
        rule.fill.fore_color.rgb = ACCENT
        rule.line.fill.background()
        rule.shadow.inherit = False

        box = sl.shapes.add_textbox(Inches(0.85), Inches(6.92),
                                    Inches(11.9), Inches(0.52))
        tf = box.text_frame
        tf.word_wrap = True

        asset = s.get("asset") or ""
        asset = asset.split("/")[-1] if "/" in asset else asset
        kind = s["kind"]
        if kind == "footage":
            kind = "ARCHIVE" if s.get("archive") else "INTERVIEW"
        head = (f"{n}.   {tc(s['prog_start'])}–{tc(s['prog_start']+s['dur'])}"
                f"   ·   {s['dur']:.2f}s   ·   {kind.upper()}"
                + (f"   ·   {asset}" if asset else ""))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = head
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = INK

        line = spoken(seg_by_i, s)
        if line:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = (line[:150] + "…") if len(line) > 150 else line
            r2.font.size = Pt(10)
            r2.font.italic = True
            r2.font.color.rgb = DIM

        notes = sl.notes_slide.notes_text_frame
        notes.text = (
            f"shot {s['name']}  |  EDL segment {s['span']}  |  "
            f"kind {s['kind']}\n"
            f"asset: {s.get('asset') or '(rendered graphic)'}\n"
            f"spec: {s['spec']}\n\n"
            f"why this picture: {s.get('why', '')}")

        if n % 40 == 0:
            print(f"[deck] {n}/{len(shots)}", flush=True)

    prs.save(str(OUT))
    mb = OUT.stat().st_size / 1048576
    print(f"[OK] {OUT}  {len(shots)} slides  {mb:.1f} MB")
    print("     slide N = shot N (no title slide)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
