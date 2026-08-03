#!/usr/bin/env python3
"""Storyboard deck: one slide per 30 seconds of programme.

Twenty-five slides instead of a hundred and seventy-two. Each slide holds
every shot in that half-minute as a labelled strip, plus the narration spoken
across it, so a whole act can be judged in one view and direction can be given
per block rather than per shot.

Shot numbers stay the ones from the 172-slide deck, so a note here
("block 6, shot 44") still points at the same thing.
"""

from __future__ import annotations

import json
import os
import subprocess
import sys
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SUFFIX = os.environ.get("BUILD_SUFFIX", "V8")
PIECES = ROOT / f"work/picture_{SUFFIX.lower()}"
STILLS = ROOT / f"work/slides_{SUFFIX.lower()}"
SHOTS = ROOT / f"manifest/picture_{SUFFIX.lower()}_shots.json"
EDL = ROOT / "manifest/edl_full.json"
WORK = ROOT / f"work/storyboard30_{SUFFIX.lower()}"
OUT = ROOT / f"final_video/MRBEAST_{SUFFIX}_STORYBOARD_30S.pptx"
FONT = "graphics/public/fonts/Anton-Regular.ttf"

BLOCK = 30.0
COLS, ROWS = 4, 2
TW, TH = 480, 270

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
GROUND = RGBColor(0x0A, 0x0A, 0x0C)
INK = RGBColor(0xEA, 0xEA, 0xEA)
DIM = RGBColor(0x9A, 0x9A, 0x9A)
ACCENT = RGBColor(0xE3, 0x12, 0x0B)


def tc(t: float) -> str:
    return f"{int(t//60)}:{t%60:05.2f}"


def mmss(t: float) -> str:
    return f"{int(t//60):02d}:{int(t%60):02d}"


def still(s: dict) -> Path:
    p = STILLS / f"{s['name']}.jpg"
    if p.exists():
        return p
    STILLS.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        ["ffmpeg", "-hide_banner", "-loglevel", "error", "-y", "-ss",
         f"{s['dur']*0.5:.3f}", "-i", str(PIECES / f"{s['name']}.mp4"),
         "-frames:v", "1", "-vf", "scale=1280:-2", "-q:v", "4", str(p)],
        check=True, timeout=300)
    return p


def esc(v: str) -> str:
    return (v.replace("\\", r"\\").replace(":", r"\:")
             .replace("'", "’").replace("%", r"\%"))


def strip(idx: int, group: list[dict]) -> Path:
    """One labelled contact strip for a 30s block."""
    dest = WORK / f"block_{idx:02d}.jpg"
    if dest.exists():
        return dest
    ins, sc = [], []
    for i, s in enumerate(group):
        ins += ["-i", str(still(s))]
        kind = s["kind"]
        if kind == "footage":
            kind = "archive" if s.get("archive") else "interview"
        lab = f"{s['_n']}  {tc(s['prog_start'])}  {kind}"
        sc.append(
            f"[{i}:v]scale={TW}:{TH},"
            f"drawbox=x=0:y={TH-30}:w={TW}:h=30:color=black@0.82:t=fill,"
            f"drawbox=x=0:y={TH-30}:w=5:h=30:color=#E3120B:t=fill,"
            f"drawtext=fontfile='{FONT}':text='{esc(lab)}'"
            f":fontcolor=#FFE04D:fontsize=21:x=11:y={TH-26}[a{i}]")
    n = len(group)
    lay = "|".join(f"{(i % COLS)*TW}_{(i//COLS)*TH}" for i in range(n))
    lb = "".join(f"[a{i}]" for i in range(n))
    filt = f"{';'.join(sc)};{lb}"
    filt += (f"xstack=inputs={n}:layout={lay}:fill=black"
             if n > 1 else "null")
    subprocess.run(
        ["ffmpeg", "-hide_banner", "-loglevel", "error", "-y", *ins,
         "-filter_complex", filt, "-frames:v", "1", "-q:v", "3", str(dest)],
        check=True, timeout=900)
    return dest


def narration(edl: dict, a: float, b: float) -> list[str]:
    out = []
    for seg in edl["segs"]:
        if seg["end"] <= a or seg["start"] >= b:
            continue
        txt = (seg.get("text") or "").strip()
        if seg["kind"] == "beat":
            txt = "[beat — music only]"
        elif seg["kind"] == "bite":
            txt = f"BITE — “{txt}”"
        if txt:
            out.append(f"{tc(seg['start'])}  {txt}")
    return out


def main() -> int:
    WORK.mkdir(parents=True, exist_ok=True)
    shots = json.loads(SHOTS.read_text(encoding="utf-8"))
    edl = json.loads(EDL.read_text(encoding="utf-8"))
    for n, s in enumerate(shots, 1):
        s["_n"] = n
    total = max(s["prog_start"] + s["dur"] for s in shots)

    blocks = []
    t = 0.0
    while t < total - 0.01:
        b = min(t + BLOCK, total)
        grp = [s for s in shots
               if s["prog_start"] < b - 0.01
               and s["prog_start"] + s["dur"] > t + 0.01]
        blocks.append((t, b, grp))
        t = b
    print(f"[storyboard] {len(blocks)} blocks of {BLOCK:.0f}s over "
          f"{total:.1f}s", flush=True)

    with ThreadPoolExecutor(max_workers=8) as ex:
        futs = [ex.submit(strip, i, g[:COLS * ROWS])
                for i, (_, _, g) in enumerate(blocks, 1) if g]
        for f in as_completed(futs):
            f.result()

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    for i, (a, b, grp) in enumerate(blocks, 1):
        sl = prs.slides.add_slide(blank)
        bg = sl.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = GROUND
        bg.line.fill.background()
        bg.shadow.inherit = False

        hd = sl.shapes.add_textbox(Inches(0.55), Inches(0.16),
                                   Inches(12.2), Inches(0.5))
        p = hd.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = (f"BLOCK {i} of {len(blocks)}      {mmss(a)} – {mmss(b)}"
                  f"      {len(grp)} shot{'s' if len(grp) != 1 else ''}"
                  + (f"      shots {grp[0]['_n']}–{grp[-1]['_n']}"
                     if grp else ""))
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = INK

        if grp:
            sl.shapes.add_picture(str(strip(i, grp[:COLS * ROWS])),
                                  Inches(0.55), Inches(0.72),
                                  Inches(12.2), Inches(3.43))
            if len(grp) > COLS * ROWS:
                warn = sl.shapes.add_textbox(Inches(0.55), Inches(4.14),
                                             Inches(12.2), Inches(0.3))
                wr = warn.text_frame.paragraphs[0].add_run()
                wr.text = (f"(+{len(grp)-COLS*ROWS} more shots in this block "
                           f"not shown)")
                wr.font.size = Pt(9)
                wr.font.color.rgb = ACCENT

        tb = sl.shapes.add_textbox(Inches(0.55), Inches(4.45),
                                   Inches(12.2), Inches(2.85))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for line in narration(edl, a, b):
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = para.add_run()
            run.text = line
            run.font.size = Pt(11)
            run.font.color.rgb = (
                INK if "BITE" in line else DIM)
            para.space_after = Pt(3)

        notes = sl.notes_slide.notes_text_frame
        notes.text = "\n".join(
            f"shot {s['_n']}  {tc(s['prog_start'])}  {s['dur']:.2f}s  "
            f"{s['kind']}  {s.get('asset') or '(rendered graphic)'}"
            for s in grp)

    prs.save(str(OUT))
    print(f"[OK] {OUT}  {len(blocks)} slides  "
          f"{OUT.stat().st_size/1048576:.1f} MB")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
